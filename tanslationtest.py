from pptx import Presentation

prs = Presentation("po.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = "shit"

prs.save("po.pptx")
print("NOW HERE")

for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            print(shape.text)
        if hasattr(shape, 'image'):
            print("Image data")
        if hasattr(shape, 'chart'):
            print("Chart data")
        if hasattr(shape, 'table'):
            print("Table data")