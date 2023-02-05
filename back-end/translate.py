import json
from flask import Flask, render_template, url_for, request, redirect, session, jsonify
from pptx import Presentation
from googletrans import Translator
import pytesseract
import imageTrans
import requests
import tanslation
api_key = 'AIzaSyA8jT11NKVpqK3HcP2ed5n30iQ3x6poffQ'
prs = Presentation("po.pptx")
#lang = 'es'

app = Flask(__name__)

#@app.route("/fetch_image", methods=["POST"])
#def fetch_image():
    # Retrieve the image from the request
 #   image = request.files["image"].read()

@app.route("/") 
def default(): 
    return render_template('setting.html')
 
@app.route("/translate/") 
def translate(): 
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:

                        text = run.text

                        # target= is the language
                        url = f'https://translation.googleapis.com/language/translate/v2?q=%7Btext%7D&target=%7Blang%7D&format=text&key=%7Bapi_key%7D'
                        response = requests.get(url)
                        if response.status_code == 200:
                            translated_text = response.json()['data']['translations'][0]['translatedText']
                            print(translated_text)
                            run.text = translated_text
                        else:
                            print('Translation failed:', response.text)
    return render_template("setting.html")


@app.route("/test/",methods = ["GET","POST"]) 
def test(): 
    if request.method == "POST":
        lang = request.form["language"]
        print(lang)
        tanslation.translate1(lang)
    return render_template("setting.html")

if __name__ == "__main__":
    app.run(debug=True)
