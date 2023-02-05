from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
from googletrans import Translator
import pytesseract
import imageTrans
import requests
import os

api_key = 'AIzaSyA8jT11NKVpqK3HcP2ed5n30iQ3x6poffQ'


def translate1(lang):

    with open("C:\\Users\\kenny\\Downloads\\filepath.txt", "r") as file:
        content = file.read()
        print(content)

    if os.path.exists("C:\\Users\\kenny\\Downloads\\filepath.txt"):
       os.remove("C:\\Users\\kenny\\Downloads\\filepath.txt")


    prs = Presentation(content)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print("The shape is a picture.")

                # Add the picture to the slide
                left = shape.left
                top = shape.top
                
                '''filepath= imageTrans.translateImg(lang)
                pic = slide.shapes.add_picture(
                    filepath, left, top,shape.width, shape.height
                )'''
            else:
                print("The shape is not a picture.")

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:

                        text = run.text

                        # target= is the language
                        url = f'https://translation.googleapis.com/language/translate/v2?q={text}&target={lang}&format=text&key={api_key}'
                        response = requests.get(url)
                        if response.status_code == 200:
                            translated_text = response.json()['data']['translations'][0]['translatedText']
                            print(translated_text)
                            run.text = translated_text
                        else:
                            print('Translation failed:', response.text)

    prs.save(content)

    print("NOW HERE")

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                print(shape.text)
            if hasattr(shape, 'image'):
                print("Image data")
            if hasattr(shape, 'chart'):
                print("Chart data")
            if hasattr(shape, 'table'):
                print("Table data")