from pptx import Presentation
from googletrans import Translator
import pytesseract

import requests

api_key = 'AIzaSyA8jT11NKVpqK3HcP2ed5n30iQ3x6poffQ'



prs = Presentation("po.pptx")
lang = 'zh-CN'

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:

                    text = run.text

                    # target= is the language
                    url = f'https://translation.googleapis.com/language/translate/v2?q={text}&target={lang}&format=text&key={api_key}'
                    response = requests.get(url)
                    if response.status_code == 200:
                        translated_text = response.json()['data']['translations'][0]['translatedText']
                        print(translated_text)
                        run.text = translated_text
                    else:
                        print('Translation failed:', response.text)


                    
prs.save("po.pptx")
print("NOW HERE")

for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            print(shape.text)
        if hasattr(shape, 'image'):
            print("Image data")
        if hasattr(shape, 'chart'):
            print("Chart data")
        if hasattr(shape, 'table'):
            print("Table data")